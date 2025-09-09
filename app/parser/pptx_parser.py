#!/usr/bin/env python3
"""
PowerPoint 파일 파서
PPT/PPTX 프레젠테이션 파일 지원
"""

import os
import logging
from typing import Dict, Any, List, Optional

from app.parser.base_parser import BaseFileParser, ParsedDocument

logger = logging.getLogger(__name__)


class PowerPointParser(BaseFileParser):
    """PowerPoint 파일 파서 클래스"""
    
    SUPPORTED_EXTENSIONS = ['.ppt', '.pptx']
    
    def __init__(self):
        super().__init__()
    
    def can_parse(self, file_path: str) -> bool:
        """PowerPoint 파일 파싱 가능 여부 확인"""
        return any(file_path.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS)
    
    def parse_file(self, file_path: str, **kwargs) -> ParsedDocument:
        """PowerPoint 파일을 파싱하여 텍스트 추출
        
        Args:
            file_path: PowerPoint 파일 경로
            **kwargs: 추가 파싱 옵션
                - extract_notes: 발표자 노트 포함 여부 (기본값: True)
                - extract_comments: 댓글 포함 여부 (기본값: False)
                
        Returns:
            ParsedDocument: 파싱된 문서 정보
            
        Raises:
            Exception: 파싱 실패 시
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            # 옵션 설정
            extract_notes = kwargs.get('extract_notes', True)
            extract_comments = kwargs.get('extract_comments', False)
            
            # PowerPoint 파일 텍스트 추출
            content, slide_count = self._extract_text_from_pptx(
                file_path, extract_notes, extract_comments
            )
            
            # 메타데이터 추출
            metadata = self._extract_metadata(file_path, slide_count)
            
            # 문서 정보 생성
            doc = ParsedDocument(
                file_path=file_path,
                file_name=os.path.basename(file_path),
                file_type='pptx',
                content=content,
                metadata=metadata,
                page_count=slide_count,
                file_size=os.path.getsize(file_path)
            )
            
            logger.info(f"Successfully parsed PowerPoint file: {file_path}")
            logger.info(f"Slides: {slide_count}, Content length: {len(content)} characters")
            
            return doc
            
        except Exception as e:
            error_msg = f"Failed to parse PowerPoint file {file_path}: {e}"
            logger.error(error_msg)
            raise Exception(error_msg) from e
    
    def _extract_text_from_pptx(self, file_path: str, extract_notes: bool = True, 
                               extract_comments: bool = False) -> tuple[str, int]:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            from pptx import Presentation
            
            # PowerPoint 파일 열기
            presentation = Presentation(file_path)
            
            content_parts = []
            slide_count = len(presentation.slides)
            
            # 각 슬라이드 처리
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = []
                
                # 슬라이드 제목
                slide_content.append(f"=== 슬라이드 {slide_num} ===")
                
                # 슬라이드의 모든 텍스트 추출
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    slide_content.append(slide_text)
                
                # 발표자 노트 추출 (옵션)
                if extract_notes and slide.has_notes_slide:
                    notes_text = self._extract_notes_text(slide)
                    if notes_text:
                        slide_content.append(f"[발표자 노트]\n{notes_text}")
                
                # 댓글 추출 (옵션)
                if extract_comments:
                    comments_text = self._extract_comments_text(slide)
                    if comments_text:
                        slide_content.append(f"[댓글]\n{comments_text}")
                
                if len(slide_content) > 1:  # 제목 외에 내용이 있는 경우
                    content_parts.append('\n'.join(slide_content))
            
            content = '\n\n'.join(content_parts)
            
            if not content.strip():
                logger.warning(f"No text content extracted from PowerPoint: {file_path}")
                content = f"PowerPoint 문서 ({os.path.basename(file_path)})에서 텍스트를 추출할 수 없습니다."
            
            return content, slide_count
            
        except ImportError:
            logger.error("python-pptx library not available")
            raise Exception("PowerPoint 파싱에 필요한 라이브러리가 설치되지 않았습니다.")
        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {e}")
            raise Exception(f"PowerPoint 파싱 중 오류 발생: {str(e)}")
    
    def _extract_slide_text(self, slide) -> str:
        """슬라이드에서 텍스트 추출"""
        text_parts = []
        
        try:
            # 슬라이드의 모든 shape 검사
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # 텍스트 정리
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)
                elif hasattr(shape, "text_frame"):
                    # 텍스트 프레임이 있는 경우
                    frame_text = self._extract_text_frame(shape.text_frame)
                    if frame_text:
                        text_parts.append(frame_text)
        
        except Exception as e:
            logger.warning(f"Failed to extract slide text: {e}")
        
        return '\n'.join(text_parts)
    
    def _extract_text_frame(self, text_frame) -> str:
        """텍스트 프레임에서 텍스트 추출"""
        text_parts = []
        
        try:
            for paragraph in text_frame.paragraphs:
                para_text = []
                for run in paragraph.runs:
                    if run.text:
                        para_text.append(run.text)
                
                if para_text:
                    text_parts.append(''.join(para_text))
        
        except Exception as e:
            logger.warning(f"Failed to extract text frame: {e}")
        
        return '\n'.join(text_parts)
    
    def _extract_notes_text(self, slide) -> str:
        """발표자 노트 텍스트 추출"""
        try:
            if not slide.has_notes_slide:
                return ""
            
            notes_slide = slide.notes_slide
            text_parts = []
            
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and "Click to add notes" not in text:  # 기본 플레이스홀더 제외
                        text_parts.append(text)
            
            return '\n'.join(text_parts)
        
        except Exception as e:
            logger.warning(f"Failed to extract notes: {e}")
            return ""
    
    def _extract_comments_text(self, slide) -> str:
        """댓글 텍스트 추출"""
        try:
            # python-pptx에서는 댓글 추출이 제한적이므로 기본 구현
            # 추후 필요시 확장 가능
            return ""
        
        except Exception as e:
            logger.warning(f"Failed to extract comments: {e}")
            return ""
    
    def _extract_metadata(self, file_path: str, slide_count: int) -> Dict[str, Any]:
        """PowerPoint 파일 메타데이터 추출"""
        metadata = {
            'file_type': 'pptx',
            'parser': 'pptx_parser',
            'slide_count': slide_count,
        }
        
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            
            # 문서 속성 추출
            core_props = presentation.core_properties
            
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.subject:
                metadata['subject'] = core_props.subject
            if core_props.created:
                metadata['created'] = core_props.created.isoformat()
            if core_props.modified:
                metadata['modified'] = core_props.modified.isoformat()
            
            # 슬라이드 크기 정보
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            metadata['slide_size'] = {
                'width': slide_width,
                'height': slide_height
            }
            
            # 레이아웃 정보
            layouts = []
            for slide in presentation.slides:
                layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
                layouts.append(layout_name)
            
            metadata['slide_layouts'] = layouts
            
        except Exception as e:
            logger.warning(f"Could not extract PowerPoint metadata: {e}")
            metadata['error'] = str(e)
        
        return metadata
    
    def validate_file(self, file_path: str) -> bool:
        """PowerPoint 파일 유효성 검사"""
        if not os.path.exists(file_path):
            return False
        
        if not self.can_parse(file_path):
            return False
        
        # 파일 크기 확인 (너무 크면 메모리 문제)
        file_size = os.path.getsize(file_path)
        if file_size > 200 * 1024 * 1024:  # 200MB 제한
            logger.warning(f"PowerPoint file too large: {file_size} bytes")
            return False
        
        # 실제 PowerPoint 파일인지 간단 확인
        try:
            if file_path.lower().endswith('.pptx'):
                # ZIP 파일인지 확인 (PPTX는 ZIP 기반)
                with open(file_path, 'rb') as f:
                    header = f.read(4)
                    if header == b'PK\x03\x04':  # ZIP 매직 바이트
                        return True
            elif file_path.lower().endswith('.ppt'):
                # 구형 PPT 파일 확인
                with open(file_path, 'rb') as f:
                    header = f.read(8)
                    if header.startswith(b'\xd0\xcf\x11\xe0'):  # OLE 매직 바이트
                        return True
        except:
            pass
        
        return True  # 확실하지 않으면 시도해보기
    
    def get_supported_extensions(self) -> List[str]:
        """지원하는 파일 확장자 목록 반환"""
        return self.SUPPORTED_EXTENSIONS.copy()


def create_parser() -> PowerPointParser:
    """PowerPoint 파서 인스턴스 생성"""
    return PowerPointParser()


# 테스트 함수
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) != 2:
        print("Usage: python pptx_parser.py <pptx_file_path>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    parser = create_parser()
    
    if not parser.validate_file(file_path):
        print(f"Invalid PowerPoint file: {file_path}")
        sys.exit(1)
    
    try:
        doc = parser.parse_file(file_path)
        print(f"File: {doc.file_name}")
        print(f"Type: {doc.file_type}")
        print(f"Size: {doc.file_size} bytes")
        print(f"Slides: {doc.page_count}")
        print(f"Content length: {len(doc.content)} characters")
        print(f"Metadata: {doc.metadata}")
        print("\n--- Content Preview ---")
        print(doc.content[:1000] + "..." if len(doc.content) > 1000 else doc.content)
        
    except Exception as e:
        print(f"Parsing failed: {e}")
        sys.exit(1)